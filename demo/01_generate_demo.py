"""Generate sample 'Ember & Oak' candle brand folder for MarkItDown demo."""
from pathlib import Path
import json, csv, random
from openpyxl import Workbook
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import LETTER
from reportlab.pdfgen import canvas
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib import colors
from PIL import Image, ImageDraw, ImageFont

ROOT = Path(__file__).parent
ROOT.mkdir(parents=True, exist_ok=True)

# 1. Q1 SALES REPORT (.xlsx)
wb = Workbook()
ws = wb.active
ws.title = "Q1 Sales by SKU"
ws.append(["SKU", "Product", "Units Sold", "Revenue (AUD)", "COGS", "Margin %"])
sales = [
    ("EO-LIB-220", "Library Smoke 220g",      724, 32_580, 11_584, 64.4),
    ("EO-FOR-220", "Forest Floor 220g",       512, 23_040,  8_192, 64.4),
    ("EO-SAL-220", "Saltwater 220g",          438, 19_710,  7_008, 64.4),
    ("EO-LIB-450", "Library Smoke 450g",       89,  6_675,  2_402, 64.0),
    ("EO-GIFT-3",  "Trio Gift Set",           104,  9_360,  4_160, 55.6),
    ("EO-WICK-1",  "Replacement Wick Pack",    61,    915,    244, 73.3),
]
for row in sales: ws.append(row)
ws.append([])
ws.append(["TOTAL", "", sum(r[2] for r in sales), sum(r[3] for r in sales), sum(r[4] for r in sales), ""])
ws2 = wb.create_sheet("Monthly Revenue")
ws2.append(["Month", "Orders", "Revenue (AUD)", "Repeat Customer %"])
ws2.append(["January 2026",  412, 28_412, 18.4])
ws2.append(["February 2026", 467, 31_287, 22.1])
ws2.append(["March 2026",    489, 32_581, 24.7])
wb.save(ROOT / "Q1-sales-report.xlsx")

# 2. SHOPIFY ORDERS (.csv)
products = [s[0:2] for s in sales]
with open(ROOT / "shopify-orders-march.csv", "w", newline="", encoding="utf-8") as f:
    w = csv.writer(f)
    w.writerow(["Order #", "Date", "Customer", "Email", "SKU", "Qty", "Total AUD", "City"])
    cities = ["Melbourne","Sydney","Brisbane","Perth","Adelaide","Hobart","Canberra","Auckland"]
    names = ["Sophie Tran","James Walker","Priya Shah","Liam O'Connor","Mia Park",
             "Noah Singh","Ava Murphy","Ethan Liu","Ruby Cooper","Oliver Patel",
             "Isla Brennan","Harvey Yu","Charlotte Reid","Finn Doyle","Zara Khan",
             "Hugo Bailey","Eden Foster","Leo Marshall","Asha Bell","Jack Sinclair"]
    for i, name in enumerate(names, start=1001):
        sku, prod = random.choice(products)
        qty = random.choice([1,1,1,2,2,3])
        unit = next(s[3]/s[2] for s in sales if s[0] == sku)
        total = round(unit * qty, 2)
        email = name.lower().replace(" ", ".").replace("'", "") + "@gmail.com"
        date = f"2026-03-{random.randint(1,31):02d}"
        w.writerow([f"#{i}", date, name, email, sku, qty, total, random.choice(cities)])

# 3. SHOPIFY PRODUCTS EXPORT (.json)
products_json = [
    {"sku": "EO-LIB-220", "title": "Library Smoke 220g",
     "price_aud": 45.00, "stock": 142,
     "scent_notes": ["leather", "smoked vanilla", "old paper"],
     "wax": "coconut-soy blend", "burn_hours": 50,
     "description": "Like the back room of a second-hand bookshop in late autumn."},
    {"sku": "EO-FOR-220", "title": "Forest Floor 220g",
     "price_aud": 45.00, "stock": 88,
     "scent_notes": ["moss", "wet bark", "cypress"],
     "wax": "coconut-soy blend", "burn_hours": 50,
     "description": "Damp earth after rain. Cold, green, alive."},
    {"sku": "EO-SAL-220", "title": "Saltwater 220g",
     "price_aud": 45.00, "stock": 64,
     "scent_notes": ["sea salt", "driftwood", "white musk"],
     "wax": "coconut-soy blend", "burn_hours": 50,
     "description": "Empty beach in winter. The wind, the rocks, the long walk back."},
    {"sku": "EO-GIFT-3", "title": "Trio Gift Set",
     "price_aud": 90.00, "stock": 32,
     "contents": ["Library Smoke 80g", "Forest Floor 80g", "Saltwater 80g"],
     "description": "All three signature scents in travel size. In a kraft box, hand-tied with linen twine."},
]
(ROOT / "shopify-products.json").write_text(json.dumps(products_json, indent=2), encoding="utf-8")

# 4. PRODUCT CATALOG (.pdf)
pdf_path = ROOT / "product-catalog-spring-2026.pdf"
doc = SimpleDocTemplate(str(pdf_path), pagesize=LETTER)
styles = getSampleStyleSheet()
story = []
story.append(Paragraph("<b>EMBER &amp; OAK</b>", styles["Title"]))
story.append(Paragraph("Spring 2026 Catalogue", styles["Heading2"]))
story.append(Spacer(1, 12))
story.append(Paragraph("Small-batch candles, hand-poured in Collingwood, Victoria. "
                       "Coconut-soy blend, cotton wicks, no synthetic dyes.", styles["BodyText"]))
story.append(Spacer(1, 18))
data = [["SKU", "Product", "Size", "Burn", "Price"]]
data.append(["EO-LIB-220", "Library Smoke", "220g", "50 hrs", "$45"])
data.append(["EO-FOR-220", "Forest Floor",  "220g", "50 hrs", "$45"])
data.append(["EO-SAL-220", "Saltwater",     "220g", "50 hrs", "$45"])
data.append(["EO-LIB-450", "Library Smoke", "450g", "95 hrs", "$75"])
data.append(["EO-GIFT-3",  "Trio Gift Set", "3x80g", "18 hrs ea", "$90"])
t = Table(data, colWidths=[90,150,60,70,60])
t.setStyle(TableStyle([
    ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#2a1a0f")),
    ("TEXTCOLOR", (0,0), (-1,0), colors.whitesmoke),
    ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
    ("GRID", (0,0), (-1,-1), 0.5, colors.grey),
    ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.HexColor("#fdf6ec"), colors.white]),
]))
story.append(t)
story.append(Spacer(1, 24))
story.append(Paragraph("<b>Founder's Note</b>", styles["Heading3"]))
story.append(Paragraph("This season we're leaning into quieter scents. Less performance, more memory. "
                       "Library Smoke remains our best-seller. Forest Floor is a slow burn that's "
                       "finally finding its people. — Maya Chen, Founder", styles["BodyText"]))
doc.build(story)

# 5. BRAND VOICE GUIDE (.docx)
docx_doc = Document()
docx_doc.add_heading("Ember & Oak — Brand Voice Guide", 0)
docx_doc.add_paragraph("Last updated: 2026-03-01")
docx_doc.add_heading("Who we are", 1)
docx_doc.add_paragraph("A small-batch candle brand based in Collingwood, Victoria. "
                       "We make candles for people who don't usually buy candles.")
docx_doc.add_heading("Voice rules", 1)
rules = [
    "Sensory, not flowery. Describe what you'd actually feel, not what marketers describe.",
    "Quiet confidence. We don't shout. We don't hype. We don't say 'luxurious'.",
    "Specific over generic. 'Wet bark after rain' beats 'fresh forest scent'.",
    "Australian English always. Favourite, colour, organise.",
    "Short sentences. Then longer ones. Rhythm matters.",
]
for r in rules:
    docx_doc.add_paragraph(r, style="List Bullet")
docx_doc.add_heading("Words we never use", 1)
docx_doc.add_paragraph("Luxurious. Indulge. Pamper. Treat yourself. Self-care. Vibes. "
                       "Notes of (just say what it smells like).")
docx_doc.add_heading("Example — bad vs good", 1)
docx_doc.add_paragraph("BAD: Indulge in our luxurious Library Smoke candle with rich notes of "
                       "vanilla and leather for the ultimate self-care moment.")
docx_doc.add_paragraph("GOOD: Library Smoke. Old paper, soft leather, a bit of vanilla. "
                       "Smells like the back room of a second-hand bookshop in late autumn.")
docx_doc.save(ROOT / "brand-voice-guide.docx")

# 6. MOTHER'S DAY CAMPAIGN (.pptx)
prs = Presentation()
prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
def add_slide(title, body_lines):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    tx = s.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12), Inches(1)).text_frame
    tx.text = title
    tx.paragraphs[0].runs[0].font.size = Pt(36)
    tx.paragraphs[0].runs[0].font.bold = True
    body = s.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(12), Inches(5)).text_frame
    body.word_wrap = True
    for i, line in enumerate(body_lines):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.runs[0].font.size = Pt(20)
add_slide("Mother's Day 2026 — Campaign Brief",
          ["Date: 11 May 2026", "Owner: Maya", "Goal: 250 Trio Gift Set sales in 10 days"])
add_slide("Audience",
          ["Adult children (28-45) shopping for mum",
           "Past customers who bought Library Smoke (highest LTV)",
           "Skews Melbourne / Sydney metro, AOV $80-120"])
add_slide("Hero offer",
          ["Trio Gift Set ($90) — free hand-written card",
           "Order by 7 May for delivery before Mother's Day",
           "Bonus: free 80g candle for orders over $150"])
add_slide("Channels",
          ["Email — 3 sends (10 days out, 5 days, 2 days)",
           "Instagram — 4 posts + Reels with founder voice",
           "Past-customer SMS on day 8 only — 'inner circle' framing"])
prs.save(ROOT / "mothers-day-campaign.pptx")

# 7. COMPETITOR PAGE (.html)
(ROOT / "competitor-page.html").write_text("""<!doctype html>
<html><head><title>Glow Co — Vanilla Bean Candle | Soy Wax | 320g</title></head>
<body>
<h1>Vanilla Bean Bliss Candle</h1>
<p class="price">$89.00 AUD</p>
<h2>Description</h2>
<p>Indulge in our luxurious Vanilla Bean Bliss candle. Hand-poured with premium soy wax,
this 320g candle features rich notes of Madagascan vanilla, warm caramel, and soft sandalwood.
Perfect for self-care moments, romantic evenings, or as a thoughtful gift. Burn time: 60 hours.</p>
<h2>Reviews (4.6/5)</h2>
<ul>
<li>"Smells amazing! So luxurious!" — Sarah M.</li>
<li>"My new favourite candle. Fills the whole room." — Jess K.</li>
<li>"Bought as a gift. She loved it." — Tom R.</li>
<li>"Burns evenly, no tunnelling. Pricey but worth it." — Anna L.</li>
</ul>
<h2>Specifications</h2>
<table>
<tr><td>Wax</td><td>100% soy</td></tr>
<tr><td>Wick</td><td>Cotton, lead-free</td></tr>
<tr><td>Burn time</td><td>60 hours</td></tr>
<tr><td>Size</td><td>320g</td></tr>
</table>
</body></html>
""", encoding="utf-8")

# 8. SUPPORT TICKETS (.csv)
with open(ROOT / "support-tickets-march.csv", "w", newline="", encoding="utf-8") as f:
    w = csv.writer(f)
    w.writerow(["Ticket #", "Date", "Customer", "SKU", "Issue", "Resolution", "Status"])
    tickets = [
        ("T-301", "2026-03-03", "Sophie Tran",     "EO-LIB-220", "Wick won't stay lit after 2 hours",         "Replacement sent",      "Resolved"),
        ("T-302", "2026-03-05", "Mia Park",        "EO-GIFT-3",  "One candle arrived broken",                   "Replacement + free wick", "Resolved"),
        ("T-303", "2026-03-07", "James Walker",    "EO-FOR-220", "Smell weaker than expected",                  "Refund offered",        "Resolved"),
        ("T-304", "2026-03-09", "Priya Shah",      "EO-LIB-220", "Wick tunneling, wax wasted",                  "Sent burn-care guide",  "Resolved"),
        ("T-305", "2026-03-12", "Oliver Patel",    "EO-SAL-220", "Question about scent strength",               "Replied with FAQ",      "Resolved"),
        ("T-306", "2026-03-14", "Charlotte Reid",  "EO-LIB-220", "Wick won't stay lit (similar to T-301)",      "Replacement sent",      "Resolved"),
        ("T-307", "2026-03-18", "Finn Doyle",      "EO-FOR-220", "Wick falling over after first burn",          "Replacement sent",      "Resolved"),
        ("T-308", "2026-03-21", "Hugo Bailey",     "EO-LIB-220", "Wick keeps drowning in wax",                  "Replacement + apology", "Resolved"),
        ("T-309", "2026-03-23", "Asha Bell",       "EO-GIFT-3",  "Gift box arrived damaged",                    "Replacement",           "Resolved"),
        ("T-310", "2026-03-26", "Jack Sinclair",   "EO-LIB-220", "Wick issue (possibly batch-related)",         "Investigating supplier","Open"),
        ("T-311", "2026-03-28", "Ruby Cooper",     "EO-SAL-220", "Loved it, just feedback — too subtle",        "Thanked, logged",       "Resolved"),
        ("T-312", "2026-03-30", "Ethan Liu",       "EO-LIB-220", "Same wick problem — fifth report this month", "Investigating supplier","Open"),
    ]
    for t in tickets: w.writerow(t)

# 9. PRODUCT MOCKUP IMAGE (.png) — text on image, for OCR demo
img = Image.new("RGB", (1200, 1200), color="#fdf6ec")
draw = ImageDraw.Draw(img)
try:
    font_big = ImageFont.truetype("arial.ttf", 90)
    font_med = ImageFont.truetype("arial.ttf", 50)
    font_sm  = ImageFont.truetype("arial.ttf", 36)
except Exception:
    font_big = ImageFont.load_default()
    font_med = font_big
    font_sm  = font_big
draw.rectangle([100, 100, 1100, 1100], outline="#2a1a0f", width=4)
draw.text((600, 280), "EMBER & OAK",        fill="#2a1a0f", font=font_big, anchor="mm")
draw.text((600, 420), "Library Smoke",      fill="#5a3820", font=font_med, anchor="mm")
draw.text((600, 510), "220g · 50 hour burn", fill="#5a3820", font=font_sm,  anchor="mm")
draw.text((600, 720), "Old paper.",          fill="#2a1a0f", font=font_med, anchor="mm")
draw.text((600, 800), "Soft leather.",       fill="#2a1a0f", font=font_med, anchor="mm")
draw.text((600, 880), "A bit of vanilla.",   fill="#2a1a0f", font=font_med, anchor="mm")
draw.text((600, 1020), "Hand-poured in Collingwood, Victoria", fill="#5a3820", font=font_sm, anchor="mm")
img.save(ROOT / "product-mockup-library-smoke.png")

# 10. FOUNDER VOICE MEMO (.txt placeholder — real version would be .mp3)
(ROOT / "founder-voice-memo-NOTE.txt").write_text(
    "[Placeholder note]\n\n"
    "In a real client folder this would be founder-voice-memo.mp3 — Maya dictating\n"
    "ideas during her morning walk. MarkItDown supports audio transcription via the\n"
    "[audio-transcription] extra. We've left this as a text placeholder so the demo\n"
    "doesn't depend on TTS or downloading a sample audio file.\n",
    encoding="utf-8"
)

# 11. README
(ROOT / "README.md").write_text("""# Ember & Oak — Sample Client Folder

Fictional small ecomm brand (artisan candles) used to demo MarkItDown's range.

## What's in here

| File | Type | Showcases |
|---|---|---|
| Q1-sales-report.xlsx | Excel | Multi-sheet workbook, tables |
| shopify-orders-march.csv | CSV | Tabular text |
| shopify-products.json | JSON | Structured data |
| product-catalog-spring-2026.pdf | PDF | Layout + tables in PDF |
| brand-voice-guide.docx | Word | Headings, bullet structure |
| mothers-day-campaign.pptx | PowerPoint | Multi-slide deck |
| competitor-page.html | HTML | Saved web page |
| support-tickets-march.csv | CSV | Patterns hidden across rows |
| product-mockup-library-smoke.png | Image | OCR (text rendered into image) |
| founder-voice-memo-NOTE.txt | (placeholder) | Audio support note |

## Demo prompts to try after running MarkItDown

- "What's our best-selling SKU and what makes it work?" (cross-references sales + products + reviews)
- "Are we losing margin anywhere I should worry about?" (sales sheet)
- "What recurring complaint shows up in support tickets?" (CSV — answer: wick problems)
- "Draft the Mother's Day email matching our brand voice." (voice guide + campaign deck)
- "How does our product copy compare to our competitor's tone?" (voice guide vs HTML)
""", encoding="utf-8")

print(f"Generated demo folder at: {ROOT}")
print(f"Files: {sorted(p.name for p in ROOT.iterdir() if p.is_file())}")
