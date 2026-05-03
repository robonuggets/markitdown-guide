"""Bulk generator — adds ~27 more office files to push total to 50.
Stays in Ember & Oak universe, content kept tight."""
from pathlib import Path
from openpyxl import Workbook
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from email.message import EmailMessage
from email import policy
from reportlab.lib.pagesizes import LETTER
from reportlab.pdfgen import canvas
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib import colors

ROOT = Path(__file__).parent

def doc(title, paragraphs, fname):
    d = Document()
    d.add_heading(title, 0)
    for p in paragraphs:
        if isinstance(p, tuple):
            kind, txt = p
            if kind == "h1": d.add_heading(txt, 1)
            elif kind == "h2": d.add_heading(txt, 2)
            elif kind == "bullet": d.add_paragraph(txt, style="List Bullet")
            else: d.add_paragraph(txt)
        else:
            d.add_paragraph(p)
    d.save(ROOT / fname)

def sheet(fname, sheets):
    wb = Workbook()
    wb.remove(wb.active)
    for name, rows in sheets:
        ws = wb.create_sheet(name)
        for r in rows: ws.append(r)
    wb.save(ROOT / fname)

def deck(fname, slides):
    prs = Presentation()
    prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
    for title, lines in slides:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        tx = s.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12), Inches(1)).text_frame
        tx.text = title
        tx.paragraphs[0].runs[0].font.size = Pt(36); tx.paragraphs[0].runs[0].font.bold = True
        body = s.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(12), Inches(5)).text_frame
        body.word_wrap = True
        for i, ln in enumerate(lines):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = ln; p.runs[0].font.size = Pt(20)
    prs.save(ROOT / fname)

def email(fname, sender, to, subj, date, body):
    m = EmailMessage(policy=policy.default)
    m["From"]=sender; m["To"]=to; m["Subject"]=subj; m["Date"]=date
    m.set_content(body)
    (ROOT / fname).write_text(str(m), encoding="utf-8")

def pdf_simple(fname, title, paragraphs):
    d = SimpleDocTemplate(str(ROOT / fname), pagesize=LETTER)
    s = getSampleStyleSheet(); story = [Paragraph(f"<b>{title}</b>", s["Title"]), Spacer(1,12)]
    for p in paragraphs:
        story.append(Paragraph(p, s["BodyText"])); story.append(Spacer(1,8))
    d.build(story)

# ============ DOCX (8) ============

doc("Returns & Refunds Policy",
    [("h1","When we refund"),
     ("bullet","Damaged on arrival — full refund or replacement, your choice."),
     ("bullet","Wick failure within 25% of advertised burn time — replacement only."),
     ("bullet","Change of mind within 14 days, unused, in original packaging — store credit, customer pays return shipping."),
     ("h1","When we don't refund"),
     ("bullet","Custom orders (engraved trios, wedding favours)."),
     ("bullet","Sale items marked 'final sale'."),
     ("bullet","Wax pull-away — this is normal in coconut-soy and not a defect."),
     ("h1","How to claim"),
     "Email returns@emberandoak.com.au with order number and a photo. We respond within 2 business days."],
    "returns-refunds-policy.docx")

doc("Hiring Notes — Junior Pourer Role",
    [("h1","Open role"),
     "Junior pourer, 32hrs/week, $28/hr, starting 12 May 2026.",
     ("h1","Candidates shortlisted"),
     ("bullet","Aiden Chow — TAFE cert IV in cosmetic chemistry, candle hobbyist, available immediately."),
     ("bullet","Mei Lin Ho — 2yrs at Mecca, strong on detail, needs 4 weeks notice."),
     ("bullet","Tom Greaves — career changer from hospitality, no candle experience, very enthusiastic."),
     ("h1","Maya's notes"),
     "Aiden has the technical chops. Mei has the discipline. Tom has the energy. Lean Aiden first, Mei as backup. Tom for casual cover during Mother's Day if we need it."],
    "hiring-notes-junior-pourer.docx")

doc("Brand Story — Long Form",
    [("h1","Beginnings"),
     "Maya Chen poured her first candle in a share-house kitchen in Brunswick in 2023. The wax was cheap, the wick was wrong, and the smell was nothing like she'd planned. She kept going.",
     ("h1","The breakthrough"),
     "Library Smoke came from a Sunday afternoon at Brown & Bunting bookshop in Northcote. Maya bought a second-hand copy of John Berger's Ways of Seeing, sat in the back room, and decided that's what she wanted a candle to smell like.",
     ("h1","Today"),
     "Four people, one workshop in Collingwood, ~2,500 candles a month. Three signature scents that haven't changed since launch. We're proud of how slow we've gone."],
    "brand-story-long.docx")

doc("Influencer Outreach Brief — Mother's Day",
    [("h1","The ask"),
     "Send 30 trios to micro-influencers (5K–25K followers) in the lifestyle/home space.",
     ("h1","Target voice"),
     ("bullet","Quiet, considered creators — not 'unboxing' channels."),
     ("bullet","Bonus if they have a parent they post about occasionally."),
     ("bullet","Skew Melbourne and Sydney."),
     ("h1","What we send"),
     "Trio gift set + handwritten card from Maya + a one-page brand story (the long-form doc).",
     ("h1","What we ask for"),
     "One organic post or story between 1-9 May. No script. We'd rather no post than a forced one."],
    "influencer-outreach-brief.docx")

doc("Workshop Safety Protocol",
    [("h1","Fire"),
     ("bullet","Class B extinguisher behind the front door — check pin monthly."),
     ("bullet","Fire blanket above the wax melter."),
     ("bullet","Never leave melter unattended above 60°C."),
     ("h1","Burns"),
     ("bullet","Cold running water, 20 minutes, then assess."),
     ("bullet","First aid kit in the office, restocked monthly."),
     ("h1","Spills"),
     ("bullet","Hot wax — let cool, scrape, then mineral spirits."),
     ("bullet","Fragrance oil on skin — soap and water, then moisturise."),
     ("h1","Evacuation"),
     "Out the back door, meet at the gate of 16 Sackville. Sam carries the headcount sheet."],
    "workshop-safety-protocol.docx")

doc("Newsletter Draft — May Edition",
    [("h1","Subject line options"),
     ("bullet","'A candle that smells like the back room of a bookshop'"),
     ("bullet","'Mother's Day, but make it specific'"),
     ("bullet","'New: Library Smoke 450g (sold out, but read on)'"),
     ("h1","Body"),
     "Hi {first_name},",
     "Mother's Day is in 11 days. We made a thing for it: the Trio Gift Set is back, hand-tied with linen twine, $90 with free shipping anywhere in Australia and NZ.",
     "We won't tell you it's 'the perfect gift'. We'll tell you it's three small candles that smell like specific places — a bookshop, a forest after rain, an empty winter beach. If your mum likes any of those things, this is for her.",
     "Order by Sunday 7 May for delivery before the day.",
     "— Maya"],
    "newsletter-draft-may.docx")

doc("Wholesale Buyer Outreach — Auckland Homestore",
    [("h1","Background"),
     "Auckland Homestore reached out 28 March via the wholesale form. They're a 4-store group across Auckland and Wellington, sell mid-range homewares.",
     ("h1","Their ask"),
     "Initial 60-unit Silver tier order across two scents (Library Smoke + Saltwater).",
     ("h1","Our terms offered"),
     "40% off RRP, Net 30, free freight on first order. NZD invoiced via Wise.",
     ("h1","Status"),
     "Quote sent 29 March, awaiting confirmation. Jordan to follow up week of 8 April if no reply."],
    "wholesale-outreach-auckland.docx")

doc("Photoshoot Shot List — May Campaign",
    [("h1","Date and location"),
     "Wed 24 April, 9am–2pm. Studio at 86 Easey St Collingwood.",
     ("h1","Shot list"),
     ("bullet","Trio set unboxing — overhead, slow hands, natural light."),
     ("bullet","Library Smoke 220g lit, on a stack of well-worn books."),
     ("bullet","Forest Floor outdoors, on damp moss (we're bringing some)."),
     ("bullet","Saltwater on driftwood — we've got pieces from St Kilda beach."),
     ("bullet","Maya hand-pouring (for the 'about us' refresh)."),
     ("h1","Crew"),
     "Photographer: Hana Sato. Stylist: Tess Marlow. Maya for hand model."],
    "photoshoot-shot-list-may.docx")

# ============ XLSX (8) ============

sheet("ad-spend-tracker.xlsx", [
    ("Meta Ads", [
        ["Campaign","Spend","Impressions","Clicks","Purchases","Revenue","ROAS"],
        ["Library Smoke retargeting", 412, 84200, 1820, 28, 1260, 3.06],
        ["Trio set prospecting",      680, 142000, 2410, 19, 1710, 2.51],
        ["Forest Floor lookalike",    288,  62100,  890, 12,  540, 1.88],
        ["Brand awareness video",     320,  98000,  410,  3,  135, 0.42],
    ]),
    ("Google Ads", [
        ["Campaign","Spend","Clicks","CPC","Purchases","ROAS"],
        ["Branded search",        140, 412, 0.34, 22, 7.07],
        ["Coconut soy candle",    280, 318, 0.88,  9, 1.45],
        ["Mother's Day candle",   400, 442, 0.91, 14, 1.58],
    ]),
])

sheet("payroll-march-2026.xlsx", [
    ("March Payroll", [
        ["Employee","Role","Hours","Rate","Gross","Super 11.5%","Net AUD"],
        ["Maya Chen",  "Director", 0, 0, 7000, 805, 5810],
        ["Sam Lee",    "Senior Pourer", 152, 34, 5168, 594, 4290],
        ["Priya Shah", "Marketing",      152, 36, 5472, 629, 4541],
        ["Jordan Wu",  "Ops + Wholesale",152, 32, 4864, 559, 4037],
    ]),
])

sheet("customer-ltv.xlsx", [
    ("LTV by Cohort", [
        ["Cohort","Customers","Avg orders","Avg AOV","12mo LTV"],
        ["Q1 2025", 312, 1.8, 71, 128],
        ["Q2 2025", 488, 2.1, 74, 155],
        ["Q3 2025", 602, 2.3, 76, 175],
        ["Q4 2025", 791, 2.0, 78, 156],
        ["Q1 2026", 894, 1.4, 78, 109],  # truncated, in progress
    ]),
    ("Top 20 Customers", [
        ["Customer","City","Orders","Total Spend AUD"],
        ["Sophie Tran",     "Melbourne", 9, 712],
        ["Ruby Cooper",     "Sydney",    7, 615],
        ["Charlotte Reid",  "Melbourne", 8, 588],
        ["Mia Park",        "Brisbane",  6, 541],
        ["Liam O'Connor",   "Hobart",    7, 522],
    ]),
])

sheet("returns-log-q1.xlsx", [
    ("Returns Q1 2026", [
        ["Date","Order","SKU","Reason","Resolution","Cost AUD"],
        ["2026-01-14","#1042","EO-LIB-220","Wick failure","Replacement",18],
        ["2026-01-22","#1057","EO-GIFT-3", "Damaged in transit","Replacement",36],
        ["2026-02-08","#1112","EO-LIB-220","Wick failure","Replacement",18],
        ["2026-02-19","#1148","EO-FOR-220","Scent too subtle","Refund",45],
        ["2026-03-03","#1201","EO-LIB-220","Wick failure","Replacement",18],
        ["2026-03-12","#1224","EO-GIFT-3", "Box dented",    "Replacement box",6],
        ["2026-03-21","#1252","EO-LIB-220","Wick failure","Replacement",18],
        ["2026-03-28","#1278","EO-LIB-220","Wick failure","Replacement",18],
    ]),
])

sheet("instagram-analytics.xlsx", [
    ("March Posts", [
        ["Date","Post","Reach","Likes","Saves","Shares","Comments"],
        ["2026-03-02","Library Smoke flat lay",     12400, 884, 142, 38, 22],
        ["2026-03-05","Workshop reel — Maya pouring",18900,1320, 88,112, 41],
        ["2026-03-09","Forest Floor candle in moss", 9200, 612, 76, 22, 11],
        ["2026-03-14","Saltwater carousel",          7800, 488, 54, 18,  9],
        ["2026-03-19","Customer UGC repost",         6100, 410, 28, 14, 12],
        ["2026-03-24","Trio gift set teaser reel",  22400,1880,302,140, 84],
        ["2026-03-30","Behind-the-scenes wax pour",  8900, 540, 92, 38, 18],
    ]),
])

sheet("gift-card-balances.xlsx", [
    ("Outstanding Gift Cards", [
        ["Card #","Issued","Original AUD","Used AUD","Balance AUD","Expires"],
        ["GC-0118","2025-12-12",100, 45, 55,"2027-12-12"],
        ["GC-0124","2025-12-19",150,150,  0,"2027-12-19"],
        ["GC-0131","2026-01-04", 50, 50,  0,"2028-01-04"],
        ["GC-0142","2026-02-14",200,  0,200,"2028-02-14"],
        ["GC-0157","2026-03-08", 90, 45, 45,"2028-03-08"],
        ["GC-0168","2026-03-22",100,  0,100,"2028-03-22"],
    ]),
])

sheet("yoy-growth.xlsx", [
    ("YoY Comparison", [
        ["Metric","Q1 2024","Q1 2025","Q1 2026","YoY %"],
        ["Revenue AUD",       11200, 67400, 92280, "+37%"],
        ["Orders",              198,   924,  1368, "+48%"],
        ["AOV AUD",              57,    73,    78,  "+7%"],
        ["Repeat customer %",  "8%", "19%",  "25%",  "+6pp"],
        ["Returns %",        "4.2%","2.8%", "1.9%", "-0.9pp"],
    ]),
])

sheet("event-budget-mothers-day-popup.xlsx", [
    ("Pop-up Budget", [
        ["Line item","Estimate AUD","Actual AUD","Notes"],
        ["Venue (Rose St Market)", 800, 850, "Two-day stall"],
        ["Travel + accommodation", 320, 280, "Two staff"],
        ["Stock pulled (cost)",   1400,1400, "100 units"],
        ["Display + signage",      450, 410, "Reusable"],
        ["Eftpos surcharges",      120,  88, ""],
        ["TOTAL",                 3090,3028, ""],
    ]),
    ("Sales", [
        ["Day","Units","Revenue AUD"],
        ["Saturday", 62, 3120],
        ["Sunday",   48, 2390],
    ]),
])

# ============ PPTX (4) ============

deck("board-update-q1.pptx", [
    ("Board Update — Q1 2026", ["Maya Chen, Founder", "April 2026"]),
    ("Numbers", ["Revenue: $92,280 (+37% YoY)", "Orders: 1,368 (+48% YoY)",
                 "Repeat customer rate: 25% (up from 19%)", "Returns: 1.9% (down from 2.8%)"]),
    ("Wins", ["Library Smoke 450g sold out in 11 days",
              "Wholesale: 3 active accounts, 12 on wait list",
              "First press hit — Broadsheet Melbourne profile in March"]),
    ("Concerns", ["Wick supplier quality issues — switching to Beecher's in May",
                  "Capacity constrained at $130K/month",
                  "Mother's Day is the largest single push, all hands needed"]),
    ("Next quarter", ["Mother's Day campaign (250-350 trios target)",
                      "Capital raise — $750K seed", "Hire 4th pourer + ops lead"]),
])

deck("brand-refresh-2026.pptx", [
    ("Brand Refresh — Internal Review", ["Priya, April 2026"]),
    ("Why now", ["Logo hasn't changed since launch", "Wholesale buyers asked for stronger shelf identity",
                 "Spring 2026 launch is a natural moment"]),
    ("Three directions", ["A — Quiet refinement (current logo, new wordmark)",
                          "B — Editorial (typeset, magazine-feel)",
                          "C — Hand-drawn (founder script)"]),
    ("Recommendation", ["Direction A. Keep equity, sharpen execution. Roll out across packaging, web, IG."]),
    ("Cost + timeline", ["Studio Fern, $14K AUD, 6 weeks. Live by 1 July."]),
])

deck("retail-strategy-2026.pptx", [
    ("Retail Strategy — H2 2026", ["Jordan, April 2026"]),
    ("Where we are", ["3 wholesale accounts, all in AU/NZ",
                      "0 owned retail", "Pop-ups: 2 in Q1, both profitable"]),
    ("The opportunity", ["Department concession — David Jones reached out",
                         "Pop-up calendar — 6 Saturdays booked across Q3",
                         "Hotel partnership — Crown Melbourne pilot in talks"]),
    ("Risk", ["Owned retail is capital intensive, distracts from product",
              "Recommend: no permanent space until 2027"]),
])

deck("2027-planning-kickoff.pptx", [
    ("2027 Planning — Early Thoughts", ["Maya, April 2026"]),
    ("Where we want to be", ["$2.5M ARR (up from ~$1.2M trajectory)",
                              "Second workshop in Sydney",
                              "8 staff total"]),
    ("New scents on the horizon", ["A late-summer scent — fig leaf, tomato vine",
                                    "A winter scent — fireplace, smoked tea",
                                    "Maybe a collab with a Melbourne perfumer"]),
    ("What we won't do", ["Diffusers, room sprays — distracts from candle craft",
                          "Subscription model — kills surprise of new scents",
                          "Marketplaces — kills margin and brand"]),
])

# ============ PDF (3) ============

pdf_simple("warehouse-pickup-form-march.pdf", "Warehouse Pickup Form — March 2026", [
    "<b>Carrier:</b> Australia Post Business",
    "<b>Pickup window:</b> Daily 2:00pm – 3:00pm",
    "<b>Total satchels March:</b> 489",
    "<b>Damaged in transit (claimed):</b> 6",
    "<b>Lost in transit:</b> 1 (Order #1244, refunded customer, claim filed)",
    "Carrier contact: Stephen Liu, 0413 882 414, Australia Post Business",
])

pdf_simple("certificate-of-analysis-wax-batch-247.pdf", "Glasshouse Wax — COA Batch GS-247", [
    "<b>Product:</b> GS-CSV-3 Coconut-Soy Wax Blend",
    "<b>Batch:</b> 247",
    "<b>Date produced:</b> 2026-03-28",
    "<b>Date shipped:</b> 2026-04-02",
    "<b>Melt point:</b> 53.4°C (spec: 52-56°C) ✓",
    "<b>Fragrance load capacity:</b> 11.2% (spec: ≥10%) ✓",
    "<b>Paraffin contamination:</b> Not detected ✓",
    "<b>QA signed:</b> Eunji Park",
])

pdf_simple("insurance-policy-summary.pdf", "Business Insurance — Policy Summary", [
    "<b>Insurer:</b> CGU Business Insurance",
    "<b>Policy number:</b> CGU-BIZ-441829",
    "<b>Annual premium:</b> $2,840 AUD",
    "<b>Coverage period:</b> 1 July 2025 – 30 June 2026",
    "<b>Public liability:</b> $20M",
    "<b>Stock cover:</b> $80K",
    "<b>Equipment cover:</b> $40K",
    "<b>Business interruption:</b> Up to $200K, 12 months",
    "<b>Workers comp:</b> Separate policy with WorkSafe Victoria",
])

# ============ EML (3) ============

email("customer-complaint-wick.eml",
    "Hannah Brennan <hannah.b@hotmail.com>", "support@emberandoak.com.au",
    "Wick won't stay lit — third candle in a row",
    "Mon, 31 Mar 2026 09:14:22 +1100",
    """Hi,

I love your candles but this is the third Library Smoke 220g in a row where the
wick drowns in wax after about 90 minutes. Order numbers: #1198, #1224, #1252.

I've been a customer for over a year and never had this issue before March. Is
there a batch problem? I'm not asking for money back, I just want a candle that
burns properly. Happy to wait if you're sorting it.

Hannah Brennan
""")

email("press-enquiry-broadsheet.eml",
    "Tessa Whitlock <tessa@broadsheet.com.au>", "maya@emberandoak.com.au",
    "Broadsheet feature — small batch candle makers in Melbourne",
    "Wed, 19 Mar 2026 11:42:08 +1100",
    """Hi Maya,

I'm writing a piece for Broadsheet on small-batch candle and fragrance makers in
inner-north Melbourne. Loved the Library Smoke (bought one last winter, still
burning it).

Could I drop by the Collingwood workshop next Tuesday or Wednesday? I'd want
about 30 minutes with you and a few photos of the pour process. We'd publish
mid-April.

Thanks,
Tessa
Broadsheet Melbourne
""")

email("landlord-rent-review.eml",
    "Crown Property <admin@crownproperty.com.au>", "maya@emberandoak.com.au",
    "Notice of rent review — 14 Sackville St Collingwood",
    "Fri, 28 Mar 2026 16:00:00 +1100",
    """Dear Tenant,

In accordance with clause 8.2 of your commercial lease for 14 Sackville Street
Collingwood, this is formal notice of the annual rent review effective 1 July 2026.

Current rent: $3,200/month (excl GST)
Proposed rent from 1 July 2026: $3,360/month (excl GST)

Increase reflects CPI movement (5.0%) for the 12 months to March 2026.

If you wish to dispute the review please respond in writing by 30 April 2026.

Regards,
Crown Property Management
""")

# ============ MISC (4) — XML, log, JSON, YAML ============

(ROOT / "klaviyo-export.xml").write_text("""<?xml version="1.0" encoding="UTF-8"?>
<klaviyo_export>
  <flow id="welcome_series" name="Welcome Series">
    <email order="1" subject="Hi from Ember &amp; Oak" open_rate="0.62" click_rate="0.18"/>
    <email order="2" subject="The story behind Library Smoke" open_rate="0.51" click_rate="0.22"/>
    <email order="3" subject="A 10% off, just because" open_rate="0.48" click_rate="0.31"/>
  </flow>
  <flow id="abandoned_cart" name="Abandoned Cart">
    <email order="1" subject="You left something in your cart" open_rate="0.44" click_rate="0.12"/>
    <email order="2" subject="Still thinking it over?" open_rate="0.31" click_rate="0.08"/>
  </flow>
  <segment name="VIP — 3+ orders" size="184" avg_ltv_aud="312"/>
  <segment name="Lapsed — no order in 6mo" size="412" avg_ltv_aud="98"/>
</klaviyo_export>
""", encoding="utf-8")

(ROOT / "shopify-webhook-log.txt").write_text("""[2026-03-30 09:14:22] orders/create #1278 — EO-LIB-220 x1, $45 — customer: ethan.liu@gmail.com
[2026-03-30 09:42:18] orders/create #1279 — EO-GIFT-3 x1, $90 — customer: claire.hong@outlook.com
[2026-03-30 10:08:44] orders/paid    #1278 — Stripe charge ch_3OqL4xK
[2026-03-30 10:18:02] inventory/update EO-LIB-220 — stock 142 -> 141
[2026-03-30 11:33:18] orders/create #1280 — EO-FOR-220 x2, EO-WICK-1 x1, $105 — customer: nick.farmer@me.com
[2026-03-30 12:01:55] orders/cancelled #1276 — customer requested, refunded $45
[2026-03-30 13:22:09] orders/fulfilled #1278 — Australia Post tracking AP4429181AU
[2026-03-30 14:48:31] orders/create #1281 — EO-GIFT-3 x1, $90 — customer: jen.maloney@gmail.com
[2026-03-30 16:14:02] inventory/update EO-LIB-450 — stock 0 -> 0 (waitlist signup count: 47)
[2026-03-30 17:55:18] orders/paid #1280 — Stripe charge ch_3OqMcL2
""", encoding="utf-8")

(ROOT / "site-config.yaml").write_text("""# Ember & Oak — Shopify theme config
brand:
  name: Ember & Oak
  tagline: Small-batch candles, hand-poured in Collingwood
  primary_colour: "#2a1a0f"
  accent_colour: "#fdf6ec"
  font_heading: "Cormorant Garamond"
  font_body: "Inter"

shipping:
  free_threshold_aud: 80
  standard_aud: 9.95
  express_aud: 14.95
  international: false  # AU + NZ only

email:
  from: hello@emberandoak.com.au
  reply_to: maya@emberandoak.com.au
  newsletter_provider: klaviyo

features:
  reviews_enabled: true
  gift_cards_enabled: true
  subscriptions_enabled: false  # decision Apr 2026 — no subs
  wholesale_portal: true
""", encoding="utf-8")

(ROOT / "supplier-contacts.json").write_text("""[
  {"name": "Glasshouse Wax Co", "category": "Wax", "contact": "David Park", "phone": "0413 887 442", "email": "david@glasshousewax.com.au", "terms": "Net 14"},
  {"name": "Atkins & Co", "category": "Wicks (legacy)", "contact": "Phil Atkins", "phone": "03 9421 8810", "email": "orders@atkinsco.com.au", "terms": "Net 30", "status": "Phasing out — quality issues"},
  {"name": "Beecher's Wicks", "category": "Wicks (new)", "contact": "Niamh Beecher", "phone": "0427 119 882", "email": "hello@beecherswicks.com.au", "terms": "Net 30", "status": "Trial starting May 2026"},
  {"name": "Brunswick Glass", "category": "Jars", "contact": "Andie Liu", "phone": "0418 442 003", "email": "wholesale@brunswickglass.com.au", "terms": "Prepay 50%"},
  {"name": "Linen Twine Co", "category": "Packaging", "contact": "Sasha M", "phone": "0431 224 119", "email": "sasha@linentwine.com.au", "terms": "Net 14"},
  {"name": "Studio Fern", "category": "Design", "contact": "Imogen Fern", "phone": "0438 990 117", "email": "imogen@studiofern.au", "terms": "50% deposit, 50% on delivery"}
]
""", encoding="utf-8")

print(f"Bulk gen done. Total files now: {sum(1 for p in ROOT.iterdir() if p.is_file())}")
