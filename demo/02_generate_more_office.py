"""Add more office files (docx, xlsx, pptx, eml) to the Ember & Oak demo folder."""
from pathlib import Path
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Inches, Pt as PPt
from email.message import EmailMessage
from email import policy

ROOT = Path(__file__).parent

# ============================================================
# WORD (.docx) — 6 files
# ============================================================

# 1. EMPLOYEE HANDBOOK
d = Document()
d.add_heading("Ember & Oak — Team Handbook", 0)
d.add_paragraph("Last updated: 2026-04-01  ·  Owner: Maya Chen")
d.add_heading("Hours and the workshop", 1)
d.add_paragraph("Pour days are Tuesday and Thursday, 8am-2pm. Pack and ship is Wednesday and Friday. "
                "Mondays are for admin, supplier calls, and the team huddle (10am, in the front room).")
d.add_heading("How we work", 1)
for line in [
    "Quiet mornings — no music until 10am, no calls before 9.",
    "Lunch at 12:30 together when we can. Bring or share, no rule.",
    "If you break a candle pouring it, just write the SKU on the whiteboard. We're tracking, not blaming.",
    "Phones in the basket during pour windows. Safety, not productivity theatre.",
]:
    d.add_paragraph(line, style="List Bullet")
d.add_heading("Leave", 1)
d.add_paragraph("20 days annual leave plus the week between Christmas and New Year (we close). "
                "Sick leave: just text Maya, no doctor's note unless it's more than three days. "
                "Mental health days count as sick days, no need to label them.")
d.add_heading("Pay", 1)
d.add_paragraph("Paid fortnightly, Thursdays. Award is Manufacturing Modern Award (MA000010). "
                "Junior pourers start at $28/hr, senior at $34/hr. Reviews twice a year.")
d.add_heading("Discount", 1)
d.add_paragraph("50% off all products for personal use, family caps at 20%. Don't resell. "
                "If you're gifting more than five candles for an occasion, just let Maya know.")
d.save(ROOT / "team-handbook.docx")

# 2. SUPPLIER AGREEMENT
d = Document()
d.add_heading("Wax Supply Agreement", 0)
d.add_paragraph("Between: Ember & Oak Pty Ltd (ABN 47 123 456 789), Collingwood VIC")
d.add_paragraph("And: Glasshouse Wax Co Pty Ltd (ABN 91 987 654 321), Brunswick VIC")
d.add_paragraph("Effective: 1 April 2026")
d.add_heading("1. Scope", 1)
d.add_paragraph("Glasshouse will supply Ember & Oak with coconut-soy wax blend (Glasshouse formula GS-CSV-3) "
                "in 25kg bags, on a rolling monthly order basis.")
d.add_heading("2. Volume and pricing", 1)
d.add_paragraph("Minimum monthly order: 200kg. Price: $7.20/kg ex GST for orders 200-500kg, "
                "$6.80/kg ex GST for orders above 500kg. Pricing locked for 12 months.")
d.add_heading("3. Quality", 1)
d.add_paragraph("Each batch must include a Certificate of Analysis covering melt point (52-56°C), "
                "fragrance load capacity (minimum 10%), and absence of paraffin contamination. "
                "If three consecutive batches fail spec, Ember & Oak may terminate without penalty.")
d.add_heading("4. Delivery", 1)
d.add_paragraph("Glasshouse delivers to 14 Sackville St Collingwood within 5 business days of order "
                "confirmation. Damaged bags replaced free of charge with photo evidence within 48 hours.")
d.add_heading("5. Payment terms", 1)
d.add_paragraph("Net 14 days from invoice. Late payments accrue 1.5% per month after grace period.")
d.add_heading("6. Termination", 1)
d.add_paragraph("Either party may terminate with 60 days written notice. Outstanding orders honoured.")
d.add_paragraph("\nSigned: Maya Chen, Director, Ember & Oak Pty Ltd")
d.add_paragraph("Signed: David Park, Director, Glasshouse Wax Co Pty Ltd")
d.save(ROOT / "supplier-agreement-glasshouse-wax.docx")

# 3. TEAM MEETING NOTES
d = Document()
d.add_heading("Team Huddle — 15 April 2026", 0)
d.add_paragraph("Present: Maya, Sam, Priya, Jordan. Apologies: none.")
d.add_heading("Updates", 1)
for line in [
    "March revenue closed at $32,581 — best month ever, up 14% on Feb.",
    "Wick supplier issue: 4 more tickets logged this week (T-313 to T-316). Sam to call Atkins Co Monday.",
    "Mother's Day campaign launches in 19 days. Maya finalising email copy this week.",
    "New scent (Library Smoke 450g) sold out — wait list at 47 names.",
]:
    d.add_paragraph(line, style="List Bullet")
d.add_heading("Decisions", 1)
for line in [
    "Switching wick supplier from Atkins to Beecher's, trial batch starts week of 22 April.",
    "Adding a fourth scent to the spring lineup — TBC, Maya to brief by 30 April.",
    "Closing Friday 25 April for ANZAC Day, no shipping that day.",
]:
    d.add_paragraph(line, style="List Bullet")
d.add_heading("Action items", 1)
d.add_paragraph("Sam: call Atkins Monday re wick batch issues. By 21 Apr.")
d.add_paragraph("Priya: photograph Trio Gift Set for Mother's Day email. By 23 Apr.")
d.add_paragraph("Maya: draft Mother's Day email v1, share for review. By 25 Apr.")
d.add_paragraph("Jordan: update wholesale price sheet for new Stockholm boutique enquiry. By 18 Apr.")
d.add_heading("Next meeting", 1)
d.add_paragraph("22 April 10am, same room.")
d.save(ROOT / "team-meeting-notes-2026-04-15.docx")

# 4. FOUNDER LETTER
d = Document()
d.add_heading("A note for the team", 0)
d.add_paragraph("From: Maya  ·  1 April 2026")
d.add_paragraph("Q1 ended on Monday. I want to say a few things before we get into the Mother's Day push.")
d.add_paragraph("First, the numbers. We did $92,280 in revenue across Jan-March. That's up 38% year on year. "
                "More importantly, repeat customer rate climbed every month — 18%, 22%, 25%. "
                "People are coming back. That's the real signal.")
d.add_paragraph("Second, the wick problem. I know it's been frustrating to keep replying to the same complaint. "
                "It's not your fault and it's not the customer's fault. Sam's chasing the supplier, "
                "we'll have a fix in two weeks. In the meantime, keep replacing freely, no questions, no quotas.")
d.add_paragraph("Third — and this is the one I keep forgetting to say out loud — the candles you're making "
                "right now are the best they've ever been. Library Smoke has a depth I couldn't get in the "
                "first six months of the brand. That's craft. That's you four. Thank you.")
d.add_paragraph("Mother's Day will be busy. We've got 10 days and a $90 hero offer. If we hit 250 trios, "
                "everyone gets a $400 bonus on the 14 May pay run. If we hit 350, it's $600.")
d.add_paragraph("Eat lunch. Drink water. Don't burn yourselves. See you Monday.")
d.add_paragraph("— Maya")
d.save(ROOT / "founder-letter-april.docx")

# 5. PRESS RELEASE
d = Document()
d.add_heading("FOR IMMEDIATE RELEASE", 0)
d.add_paragraph("Contact: press@emberandoak.com.au  ·  Maya Chen, Founder")
d.add_heading("Ember & Oak launches Spring 2026 collection from Collingwood workshop", 1)
d.add_paragraph("MELBOURNE, 1 April 2026 — Ember & Oak, the small-batch candle brand based in Collingwood, "
                "today announced its Spring 2026 collection alongside results showing 38% year-on-year "
                "revenue growth in Q1.")
d.add_paragraph("The new lineup adds a 450g Library Smoke (the brand's best-selling scent, now in a longer-burn "
                "size) and reintroduces the Trio Gift Set for Mother's Day. All candles continue to use "
                "founder Maya Chen's signature coconut-soy blend, hand-poured in the brand's Sackville Street workshop.")
d.add_paragraph("\"We're not chasing trends,\" said Chen. \"Library Smoke smells like a second-hand bookshop, "
                "and we sell out every month. People are tired of being sold 'luxury' — they want a candle "
                "that smells like a real place.\"")
d.add_paragraph("Ember & Oak ships across Australia and New Zealand, with wholesale partnerships in Sydney, "
                "Melbourne, Brisbane and Auckland. The Spring 2026 collection is available now at "
                "emberandoak.com.au, with prices from $45.")
d.add_paragraph("\nAbout Ember & Oak: founded in 2024 by Maya Chen, a former architectural lighting designer, "
                "Ember & Oak makes coconut-soy candles in Collingwood, Victoria. The four-person team "
                "produces around 2,500 candles per month. www.emberandoak.com.au")
d.add_paragraph("\n— ENDS —")
d.save(ROOT / "press-release-spring-2026.docx")

# 6. STORE OPENING CHECKLIST
d = Document()
d.add_heading("Workshop Opening Checklist — Daily", 0)
d.add_paragraph("Owner: Sam (lead pourer). Print and tick.")
d.add_heading("Before 8am", 1)
for line in [
    "Unlock back door, alarm code 4419.",
    "Turn on extraction fan (switch left of the pour bench).",
    "Heat wax melter to 75°C, takes 25 minutes.",
    "Check temperature log book — note overnight low.",
    "Pull today's order sheet from printer tray.",
]:
    d.add_paragraph(line, style="List Bullet")
d.add_heading("Before pouring", 1)
for line in [
    "Wipe down pour bench with damp cloth.",
    "Lay out wicks for today's batch (count from order sheet, +5%).",
    "Test fragrance bottles — shake, sniff, note any off-batches.",
    "Photo of empty bench for the daily Slack channel.",
]:
    d.add_paragraph(line, style="List Bullet")
d.add_heading("End of pour day", 1)
for line in [
    "Wax melter off. Confirm by touching the cord — should be cold by 4pm.",
    "Sweep floor (broom in the storeroom corner).",
    "Empty bin into the green skip out the back.",
    "Lock back door, set alarm, confirm flashing red.",
]:
    d.add_paragraph(line, style="List Bullet")
d.save(ROOT / "workshop-opening-checklist.docx")

# ============================================================
# EXCEL (.xlsx) — 4 files
# ============================================================

# 7. Q2 BUDGET FORECAST (multi-sheet)
wb = Workbook()
ws = wb.active
ws.title = "Q2 Forecast"
ws.append(["Category", "April", "May", "June", "Q2 Total"])
budget = [
    ("Revenue (target)",        38000, 52000, 41000, 131000),
    ("Wax + fragrance",         -4200, -5800, -4500, -14500),
    ("Wicks + jars",            -3800, -5200, -4100, -13100),
    ("Packaging + shipping",    -3100, -4400, -3300, -10800),
    ("Wages (4 staff)",        -16800,-16800,-16800, -50400),
    ("Rent (Collingwood)",      -3200, -3200, -3200,  -9600),
    ("Marketing + ads",         -2500, -4000, -2000,  -8500),
    ("Software + admin",        -1200, -1200, -1200,  -3600),
    ("EBITDA",                   3200, 11400,  5900,  20500),
]
for row in budget:
    ws.append(row)
for cell in ws["1:1"]:
    cell.font = Font(bold=True, color="FFFFFF")
    cell.fill = PatternFill("solid", fgColor="2A1A0F")

ws2 = wb.create_sheet("Assumptions")
ws2.append(["Assumption", "Value", "Notes"])
ws2.append(["Avg order value", 78, "Tracking up from $71 in Q1"])
ws2.append(["Mother's Day uplift", "+38%", "Based on 2025 May vs Apr"])
ws2.append(["Wholesale share", "22%", "Stable, three accounts"])
ws2.append(["Repeat customer rate", "26%", "Q1 closed at 25%, modest growth assumed"])
ws2.append(["Ad spend ROAS target", "3.5x", "Q1 actual was 3.2x"])

ws3 = wb.create_sheet("Cash Position")
ws3.append(["Date", "Opening", "Inflow", "Outflow", "Closing"])
ws3.append(["1-Apr", 47200, 38000, -34800,  50400])
ws3.append(["1-May", 50400, 52000, -40600,  61800])
ws3.append(["1-Jun", 61800, 41000, -35100,  67700])
wb.save(ROOT / "q2-budget-forecast.xlsx")

# 8. INVENTORY SNAPSHOT
wb = Workbook()
ws = wb.active
ws.title = "Stock On Hand"
ws.append(["SKU", "Product", "Stock", "Reorder At", "Status"])
inv = [
    ("EO-LIB-220", "Library Smoke 220g",      142, 100, "OK"),
    ("EO-FOR-220", "Forest Floor 220g",        88, 100, "REORDER"),
    ("EO-SAL-220", "Saltwater 220g",           64, 100, "REORDER"),
    ("EO-LIB-450", "Library Smoke 450g",        0,  20, "OUT — 47 on waitlist"),
    ("EO-GIFT-3",  "Trio Gift Set",            32,  40, "REORDER"),
    ("EO-WICK-1",  "Replacement Wick Pack",   210, 100, "OK"),
]
for row in inv:
    ws.append(row)
ws2 = wb.create_sheet("Raw Materials")
ws2.append(["Material", "Unit", "On Hand", "Weeks Cover", "Reorder?"])
ws2.append(["Coconut-soy wax", "kg",       340, 5.2, "No"])
ws2.append(["Cotton wicks (small)", "units", 1800, 4.1, "No"])
ws2.append(["Cotton wicks (large)", "units",  240, 1.8, "Yes — order 1000"])
ws2.append(["Library Smoke fragrance", "kg",   12, 3.0, "Yes — order 10kg"])
ws2.append(["Forest Floor fragrance", "kg",    18, 4.5, "No"])
ws2.append(["Saltwater fragrance", "kg",        8, 2.6, "Yes — order 8kg"])
ws2.append(["220g amber jars", "units",       620, 3.9, "Yes — order 500"])
ws2.append(["450g amber jars", "units",        15, 0.4, "URGENT — order 300"])
wb.save(ROOT / "inventory-snapshot.xlsx")

# 9. WHOLESALE PRICING
wb = Workbook()
ws = wb.active
ws.title = "Wholesale Tiers"
ws.append(["Tier", "Min Order (units)", "Discount Off RRP", "Payment Terms", "Example Stockists"])
ws.append(["Bronze",  24,  "30%", "Net 30",            "Café stockists, gift shops"])
ws.append(["Silver",  60,  "40%", "Net 30",            "Boutique homewares"])
ws.append(["Gold",   120,  "45%", "Net 14, prepay 50%", "Department concessions"])
ws.append(["House",  300,  "50%", "Prepay 100%",        "Hotels, large retailers"])

ws2 = wb.create_sheet("Pricing by SKU")
ws2.append(["SKU", "RRP", "Bronze", "Silver", "Gold", "House"])
prices = [
    ("EO-LIB-220", 45),
    ("EO-FOR-220", 45),
    ("EO-SAL-220", 45),
    ("EO-LIB-450", 75),
    ("EO-GIFT-3",  90),
    ("EO-WICK-1",  15),
]
for sku, rrp in prices:
    ws2.append([sku, rrp, round(rrp*0.7,2), round(rrp*0.6,2), round(rrp*0.55,2), round(rrp*0.5,2)])

ws3 = wb.create_sheet("Active Accounts")
ws3.append(["Account", "City", "Tier", "Last Order", "Next Order Due"])
ws3.append(["The Tea Room",          "Sydney",    "Bronze", "2026-03-12", "2026-04-15"])
ws3.append(["Pickett & Hide",        "Melbourne", "Silver", "2026-03-04", "2026-04-08"])
ws3.append(["Auckland Homestore",    "Auckland",  "Silver", "2026-02-28", "2026-04-20"])
wb.save(ROOT / "wholesale-pricing.xlsx")

# 10. EXPENSE REPORT
wb = Workbook()
ws = wb.active
ws.title = "March Expenses"
ws.append(["Date", "Vendor", "Category", "Description", "Amount AUD", "GST Inc"])
exp = [
    ("2026-03-02", "Glasshouse Wax Co",  "Materials",  "200kg coconut-soy wax",          1584.00, "Y"),
    ("2026-03-04", "Atkins & Co",        "Materials",  "Cotton wicks, 2000 units",        420.00, "Y"),
    ("2026-03-05", "Australia Post",     "Shipping",   "March bulk satchels",             892.50, "Y"),
    ("2026-03-08", "Officeworks",        "Admin",      "Printer ink, paper",              112.30, "Y"),
    ("2026-03-09", "Square",             "Software",   "POS subscription",                 89.00, "Y"),
    ("2026-03-12", "Shopify",            "Software",   "Plus plan, March",                452.00, "Y"),
    ("2026-03-14", "Klaviyo",            "Software",   "Email tier, March",               180.00, "Y"),
    ("2026-03-15", "ANZ Bank",           "Fees",       "Merchant fees Feb",               341.20, "N"),
    ("2026-03-18", "Meta Ads",           "Marketing",  "Instagram + Facebook ads",       1200.00, "Y"),
    ("2026-03-19", "Google Ads",         "Marketing",  "Search campaigns",                820.00, "Y"),
    ("2026-03-21", "Coles",              "Workshop",   "Tea, coffee, milk for studio",     67.40, "Y"),
    ("2026-03-22", "Energy Australia",   "Utilities",  "Workshop electricity",            286.50, "Y"),
    ("2026-03-25", "Internet Provider",  "Utilities",  "Workshop NBN, March",             89.95, "Y"),
    ("2026-03-26", "Brunswick Glass",    "Materials",  "Amber jars 220g x 500",           895.00, "Y"),
    ("2026-03-27", "Maya Chen",          "Reimburse",  "Postage stamps, packing tape",     34.20, "Y"),
    ("2026-03-29", "MYOB",               "Software",   "Accounting subscription",          78.00, "Y"),
    ("2026-03-30", "Crown Property",     "Rent",       "Workshop rent, April",           3200.00, "Y"),
]
for row in exp:
    ws.append(row)
ws.append([])
ws.append(["", "", "", "TOTAL", sum(r[4] for r in exp), ""])
wb.save(ROOT / "expense-report-march.xlsx")

# ============================================================
# POWERPOINT (.pptx) — 2 files
# ============================================================

def add_slide(prs, title, body_lines):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    tx = s.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12), Inches(1)).text_frame
    tx.text = title
    tx.paragraphs[0].runs[0].font.size = PPt(36)
    tx.paragraphs[0].runs[0].font.bold = True
    body = s.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(12), Inches(5)).text_frame
    body.word_wrap = True
    for i, line in enumerate(body_lines):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.runs[0].font.size = PPt(20)

# 11. INVESTOR PITCH
prs = Presentation()
prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
add_slide(prs, "Ember & Oak — Series Seed Pitch",
          ["Maya Chen, Founder", "April 2026", "Raising $750K AUD on $4M post-money"])
add_slide(prs, "The candle market is full of liars",
          ["$2.1B AUD market, growing 6% per year",
           "Saturated with 'luxury' marketing and synthetic vanilla",
           "Customers are tired — average review sentiment trending down",
           "Opportunity for an honest, scent-led brand"])
add_slide(prs, "What we built",
          ["Coconut-soy candles, hand-poured in Collingwood",
           "Three signature scents — Library Smoke, Forest Floor, Saltwater",
           "Direct + wholesale, no marketplaces",
           "Q1 2026: $92K revenue, 25% repeat rate, 38% YoY growth"])
add_slide(prs, "Traction",
          ["Q1 2024: $11K", "Q1 2025: $67K", "Q1 2026: $92K",
           "Library Smoke has been our #1 SKU for 14 consecutive months"])
add_slide(prs, "Use of funds",
          ["$300K — second workshop space in Sydney",
           "$200K — hire 4 more pourers + ops lead",
           "$150K — 12 months ad budget across Meta + Google",
           "$100K — buffer + working capital"])
add_slide(prs, "Why now",
          ["At $130K monthly revenue we're capacity-constrained",
           "Wholesale wait list: 12 stockists, 4 markets we can't serve",
           "Mother's Day 2026 will be our largest single campaign — proof point"])
add_slide(prs, "The team",
          ["Maya Chen — Founder. Former architectural lighting designer.",
           "Sam Lee — Head pourer, with us since launch.",
           "Priya Shah — Marketing, ex-Aesop.",
           "Jordan Wu — Ops + wholesale."])
add_slide(prs, "Ask",
          ["$750K AUD on $4M post-money",
           "18-month runway to $5M ARR",
           "Lead investor identified — looking for two more"])
prs.save(ROOT / "investor-pitch-april-2026.pptx")

# 12. PACKAGING REDESIGN
prs = Presentation()
prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
add_slide(prs, "Packaging Redesign — Why and What",
          ["Internal review", "Maya + Priya", "April 2026"])
add_slide(prs, "The problem",
          ["Current kraft box is sturdy but generic",
           "Five customer photos this quarter showed boxes arriving dented",
           "Unboxing video views drop off in the first 4 seconds",
           "Wholesale buyers asked for 'something more giftable'"])
add_slide(prs, "What we keep",
          ["Brown kraft palette — it's us",
           "Linen twine tie",
           "Hand-stamped wax seal on Trio Gift Set"])
add_slide(prs, "What changes",
          ["Stronger 350gsm board (up from 280gsm)",
           "Internal candle cradle — moulded pulp, not foam",
           "Belly band in deep ochre with letterpress logo",
           "Inside lid — a printed 'how to burn it' card, not a sticker"])
add_slide(prs, "Cost impact",
          ["Per-unit packaging cost rises from $1.40 to $2.20",
           "On a $45 candle: 1.8% margin hit",
           "Offset by removing the fragile-sticker step (saves 8 sec/candle)"])
add_slide(prs, "Timeline",
          ["April: finalise design with Studio Fern",
           "May: prototype 50 boxes, test ship to 10 staff homes",
           "June: full transition to new packaging",
           "Old stock burned through during Mother's Day"])
prs.save(ROOT / "packaging-redesign-rationale.pptx")

# ============================================================
# EMAIL (.eml) — 1 file
# ============================================================

# 13. SUPPLIER EMAIL (RFC 822)
msg = EmailMessage(policy=policy.default)
msg["From"] = "David Park <david@glasshousewax.com.au>"
msg["To"] = "Maya Chen <maya@emberandoak.com.au>"
msg["Cc"] = "Sam Lee <sam@emberandoak.com.au>"
msg["Subject"] = "Re: April wax order + heads-up on May pricing"
msg["Date"] = "Wed, 16 Apr 2026 14:22:18 +1000"
msg.set_content("""Hi Maya,

Confirming your April order — 250kg of GS-CSV-3, locked at $7.20/kg ex GST per the
agreement we signed on the 1st. Ship date Tuesday 22 April. Sam, I'll text you
when the truck leaves Brunswick.

Two things to flag for May.

1. Coconut wax is up about 8% from our supplier in the Philippines (typhoon season
   hit harder than expected). I'll absorb that for April since we've locked
   pricing, but if it stays elevated through Q3 we may need to re-open the
   contract for the back half of the year. I'll send numbers in late May.

2. We're trialling a new soy stream from a regen-ag farm in NSW. If your customers
   care about that story, I can send a sample bag with the April delivery — you
   could pour a small batch and see if you notice the difference. Up to you.

Also — saw the Library Smoke 450g sold out. Congrats. Let me know if you need
priority on May volume to refill that line.

Cheers
Dave

David Park
Director, Glasshouse Wax Co
0413 887 442
""")
(ROOT / "supplier-email-glasshouse.eml").write_text(str(msg), encoding="utf-8")

# ============================================================
# UPDATE README
# ============================================================
(ROOT / "README.md").write_text("""# Ember & Oak — Sample Client Folder

Fictional small ecomm brand (artisan candles) used to demo MarkItDown's range.
23 source files across 9 formats. Run `_run_markitdown.py` to flatten everything to `_markdown/`.

## What's in here

### Original demo files

| File | Type | Showcases |
|---|---|---|
| Q1-sales-report.xlsx | Excel | Multi-sheet workbook, tables |
| shopify-orders-march.csv | CSV | Tabular text |
| shopify-products.json | JSON | Structured data |
| product-catalog-spring-2026.pdf | PDF | Layout + tables in PDF |
| brand-voice-guide.docx | Word | Headings, bullet structure |
| mothers-day-campaign.pptx | PowerPoint | Multi-slide deck |
| competitor-page.html | HTML | Saved web page |
| support-tickets-march.csv | CSV | Patterns hidden across rows |
| product-mockup-library-smoke.png | Image | OCR (text rendered into image) |
| founder-voice-memo-NOTE.txt | Plain text | Audio placeholder |

### More office files (added April 2026)

| File | Type | Showcases |
|---|---|---|
| team-handbook.docx | Word | Internal HR doc with sections |
| supplier-agreement-glasshouse-wax.docx | Word | Contract style with numbered clauses |
| team-meeting-notes-2026-04-15.docx | Word | Meeting minutes + action items |
| founder-letter-april.docx | Word | Long-form letter, no headings |
| press-release-spring-2026.docx | Word | Press release format |
| workshop-opening-checklist.docx | Word | Operational checklist with bullets |
| q2-budget-forecast.xlsx | Excel | 3 sheets — forecast, assumptions, cash |
| inventory-snapshot.xlsx | Excel | Stock on hand + raw materials |
| wholesale-pricing.xlsx | Excel | Tiered pricing + active accounts |
| expense-report-march.xlsx | Excel | Line-item ledger |
| investor-pitch-april-2026.pptx | PowerPoint | 8-slide fundraising deck |
| packaging-redesign-rationale.pptx | PowerPoint | Internal design review deck |
| supplier-email-glasshouse.eml | Email | Outlook/RFC 822 message |

## Demo prompts to try after running MarkItDown

- "What's our best-selling SKU and what makes it work?"
- "Are we losing margin anywhere I should worry about?" (sales + expense report)
- "What recurring complaint shows up in support tickets?"
- "Draft the Mother's Day email matching our brand voice."
- "Summarise the supplier agreement in plain English."
- "What's our cash position heading into June?" (Q2 forecast)
- "Which raw materials need reordering urgently?"
- "Pull the action items from last week's huddle."
- "Did the supplier flag any pricing risks for May?" (eml)
- "Compare our investor pitch numbers against the Q2 forecast — do they line up?"
""", encoding="utf-8")

print(f"Generated additional files at: {ROOT}")
print(f"Total files now: {sum(1 for p in ROOT.iterdir() if p.is_file())}")
print("New files:")
for p in sorted(ROOT.iterdir()):
    if p.is_file() and p.name not in {"README.md", "_generate_demo.py", "_run_markitdown.py", "_ocr_image.py", "_generate_more_office.py"}:
        print(f"  {p.name}")
